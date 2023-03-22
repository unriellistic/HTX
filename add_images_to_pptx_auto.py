"""
The script automates the creation of multiple slides with a specific content:
- One picture
- Two subtexts (Filename, Remarks)

Update the different directories.

EYEFOX_IMG_DIR: Is the full absolute path to the images taken by the eyefox machine
XRAY_IMG_DIR: Is the full absolute path to the images taken by the xray machine. (Tiff format ones)
EXCEL_FILE_DIR = Is the full absolute path to the stats noted down on the different objects.
    - See below for how to format the excel file, and alter the code itself to take in the correct column header information.
EXCEL_FILE_SHEETNAME: Name of excel file sheetname that you take the stats from.
PPTX_FILENAME: Name of pptx file that you want to save as.
"""

from pptx import Presentation
import os

### ------------------------- Constants to change ----------------------------------
# EYEFOX_IMG_DIR: folder which stores eyefox images
# XRAY_IMG_DIR: folder which stores Xray images
# EXCEL_FILE_DIR: path to the excel file that stores the info
# EXCEL_FILE_SHEETNAME: name of excel sheet, default is "Sheet1"
# PPTX_FILENAME = name of pptx file to save as

EYEFOX_IMG_DIR = "D:\\Eyefox Data\\CAG\\Staged Scene\\16 - Cleaned Images Annotated"
XRAY_IMG_DIR = "D:\\Eyefox Data\\CAG\\Staged Scene\\16 - Cleaned TIF X-Ray Images"
EXCEL_FILE_DIR = "D:\\Eyefox Data\\CAG\\Staged Scene\\EYEFOX_STATS_16JAN.xls"
EXCEL_FILE_SHEETNAME = "Data for 160123"
PPTX_FILENAME = "test"

### --------------------------------------------------------------------------------

### ------------------------ Slide class -----------------------------------------
prs = Presentation()
# If want to see slide layout, uncomment below
# for slide in prs.slide_layouts:
#     for shape in slide.placeholders:
#         print('%d %d %s' % (prs.slide_layouts.index(slide), shape.placeholder_format.idx, shape.name))
class MySlide:
    def __init__(self, data):
        self.layout = prs.slide_layouts[data[3]]
        self.slide = prs.slides.add_slide(self.layout)
        self.title = self.slide.shapes.title
        self.title.text = data[0]
        self.subtitle = self.slide.placeholders[1]
        self.subtitle.text = data[1]

        # Uncomment below if want to see placeholder information
        # for shape in self.slide.placeholders:
        #     print('%d %s %s' % (
        #         shape.placeholder_format.idx,
        #         shape.placeholder_format.type,
        #         shape.name))
        #     print()
        if data[2] != "":
            self.img = self.slide.placeholders[1].insert_picture(data[2])
            available_width = self.img.width
            available_height = self.img.height
            image_width, image_height = self.img.image.size
            placeholder_aspect_ratio = float(available_width) / float(available_height)
            image_aspect_ratio = float(image_width) / float(image_height)

            # Get initial image placeholder left and top positions
            pos_left, pos_top = self.img.left, self.img.top

            self.img.crop_top = 0
            self.img.crop_left = 0
            self.img.crop_bottom = 0
            self.img.crop_right = 0

            # ---if the placeholder is "wider" in aspect, shrink the self.img width while
            # ---maintaining the image aspect ratio
            if placeholder_aspect_ratio > image_aspect_ratio:
                self.img.width = int(image_aspect_ratio * available_height)
                self.img.height = available_height

            # ---otherwise shrink the height
            else:
                self.img.height = int(available_width / image_aspect_ratio)
                self.img.width = available_width

            # Set the self.img left and top position to the initial placeholder one
            self.img.left, self.img.top = pos_left, pos_top

            # Or if we want to center it vertically:
            # self.img.top = self.img.top + int(self.img.height/2)
### ---------------------------------------------------------------------------

### --------------------- Get information from excel ------------------------
import pandas as pd
df = pd.read_excel(EXCEL_FILE_DIR, sheet_name=EXCEL_FILE_SHEETNAME)

print("Excel columns:", df.columns.tolist())
info_about_images = []

for index, row in df.iterrows():
    # Overall = Whether Eyefox got detect
    # Filename = Top Filename (View 1)
    # 'Unnamed: 15' = Side Filename (View 2)
    # 'Unnamed: 16' = X-Ray Filename...""
    if index == 0:
        continue
    else:
        # NOTE: If want to alter code for different excel, it's the code located 2 lines below this sentence to edit the columns to take the information from.
        # print(f"Current row: {str(index)}\n", row['Overall'], row['Top View (view 1)'], row['Side View (view 2)'], row['Filename'], row['Unnamed: 15'], row['Unnamed: 16'], row['Legend/Notes'])
        info_about_images.append(((True if row['Overall'] == 1 else False), # input column header for overall detection
                                 (True if row['Top View (view 1)']==1 else False), # input column header for top view detection
                                 (True if row['Side View (view 2)']==1 else False), # input column header for side view detection
                                 str(row['Filename']), # input eyefox top filename column header
                                 str(row['Unnamed: 15']), # input eyefox side filename column header
                                 str(row['Unnamed: 16']), # input xray filename column header
                                 str(row['Legend/Notes'])))
### --------------------------------------------------------------------
### ------------ Clean information and put into pptx -------------------
from general_scripts import change_file_extension
log = []
for set, detail in enumerate(info_about_images):
    
    # Detail is a tuple with information stored in this order:
    # 0: Overall detection
    # 1: Top view detection
    # 2: Side view detection
    # 3: Eyefox top filename
    # 4: Eyefox side filename
    # 5: XRAY filename
    # 6: Legend/notes
    eyefox_detection = ""
    top_view_present = False
    
    ### Rename json files to image files
    
    # Check if filename was blank
    if detail[4] != "nan":
        try:
            side_eyefox_image_filename = change_file_extension(detail[4], ".png")
        except:
            print(f"error in: {detail[4]}")
            log.append(detail[4])

    # Check if filename was blank
    if detail[3] != "nan":
        try:
            top_eyefox_image_filename = change_file_extension(detail[3], ".png")
        except:
             print(f"error in: {detail[3]}")
             log.append(detail[3])

    ## Create xray slides

    # If overall detection is True
    if detail[0] == True:
        eyefox_detection = eyefox_detection + " YES"
    else:
        eyefox_detection = " NO"
    top_xray_image = XRAY_IMG_DIR + "\\" + detail[5] + "_0_B.tif"
    side_xray_image = XRAY_IMG_DIR + "\\" + detail[5] + "_0_A.tif"
    
    # Placeholder information for slide layout 8
    # 8 0 Title 1
    # 8 1 Picture Placeholder 2
    # 8 2 Text Placeholder 3
    # 8 10 Date Placeholder 4
    # 8 11 Footer Placeholder 5
    # 8 12 Slide Number Placeholder 6

    # 0 TITLE (1) Title 1
    # 1 PICTURE (18) Picture Placeholder 2
    # 2 BODY (2) Text Placeholder 3

    ## Create eyefox slides
    try:
        MySlide([f"File: {detail[5]}\nEyefox detection:{eyefox_detection}\nRemarks: {detail[6]}", "", side_xray_image, 8])
        MySlide([f"File: {detail[5]}\nEyefox detection:{eyefox_detection}\nRemarks: {detail[6]}", "", top_xray_image, 8])
    except:
        print(f"Can't find: {side_xray_image}")
        # This is to account for empty x-ray file name found. Likely due to 2 rows being merged because 1 image contains multiple threats.
        # Continue so that skips everything else.
        continue

    # If Side View detection is True
    if detail[2] == True:
        side_eyefox_image = EYEFOX_IMG_DIR + "\\" + side_eyefox_image_filename
        try:
            MySlide([f"File: {side_eyefox_image_filename}\nEyefox detection:{eyefox_detection}\nRemarks: {detail[6]}", "eyefox_test", side_eyefox_image, 8])
        except:
            print(f"Can't find: {side_eyefox_image_filename}")
            log.append(side_eyefox_image_filename)

    # If Top View detection is True
    if detail[1] == True:
        top_eyefox_image = EYEFOX_IMG_DIR + "\\" + top_eyefox_image_filename
        try:
            # Have to leave an empty '""' because the pptx class is formatted in such a way.
            MySlide([f"File: {top_eyefox_image_filename}\nEyefox detection:{eyefox_detection}\nRemarks: {detail[6]}", "", top_eyefox_image, 8])
        except:
            print(f"Can't find: {top_eyefox_image_filename}")
            log.append(top_eyefox_image_filename)

print(f"Saving '{PPTX_FILENAME}.pptx'...")
prs.save("test.pptx")
print(f"Starting '{PPTX_FILENAME}.pptx'...")
os.startfile("test.pptx")

print("Error log:", log)